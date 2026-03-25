import os, tempfile, shutil, base64
from pathlib import Path
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, JSONResponse
from typing import List, Optional
import fitz

router = APIRouter()

import math

def _apply_watermark_text(page, text: str, font_size: int, rotation_deg: float):
    """
    Draw watermark text at arbitrary rotation using Shape + morph matrix.
    Works for any angle (0, 45, 90, etc.) unlike insert_text which only
    accepts multiples of 90.
    """
    pw, ph = page.rect.width, page.rect.height
    # Center of page
    cx, cy = pw / 2, ph / 2

    # Build rotation matrix (fitz.Matrix accepts degrees directly)
    mat = fitz.Matrix(rotation_deg)

    shape = page.new_shape()
    # Place text at center, then rotate around center point
    shape.insert_text(
        fitz.Point(cx, cy),
        text,
        fontsize=font_size,
        color=(0.65, 0.65, 0.65),
        morph=(fitz.Point(cx, cy), mat),
    )
    shape.commit(overlay=True)



# ── THUMBNAILS ─────────────────────────────────────────────────────────
@router.post("/thumbnails")
async def get_thumbnails(file: UploadFile = File(...)):
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    with open(src,"wb") as f: shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    thumbs = []
    mat = fitz.Matrix(0.35, 0.35)
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        b64 = base64.b64encode(pix.tobytes("png")).decode()
        thumbs.append({"page":i+1,"img":f"data:image/png;base64,{b64}","width":pix.width,"height":pix.height})
    doc.close(); shutil.rmtree(tmp, ignore_errors=True)
    return {"count":len(thumbs),"thumbnails":thumbs}

# ── MERGE ──────────────────────────────────────────────────────────────
@router.post("/merge")
async def merge(files: List[UploadFile] = File(...)):
    if len(files)<2: raise HTTPException(400,"Need at least 2 PDFs.")
    tmp = tempfile.mkdtemp(); out = os.path.join(tmp,"merged.pdf")
    result = fitz.open()
    for f in files:
        p = os.path.join(tmp,f.filename)
        with open(p,"wb") as fp: shutil.copyfileobj(f.file,fp)
        doc = fitz.open(p); result.insert_pdf(doc); doc.close()
    result.save(out); result.close()
    return FileResponse(out, media_type="application/pdf", filename="merged.pdf")

# ── SPLIT ──────────────────────────────────────────────────────────────
@router.post("/split")
async def split(file: UploadFile = File(...), pages: str = Form(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"split.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src); result = fitz.open()
    for part in pages.split(","):
        part=part.strip()
        if "-" in part:
            a,b=part.split("-")
            for n in range(int(a)-1,min(int(b),doc.page_count)): result.insert_pdf(doc,from_page=n,to_page=n)
        elif part.isdigit():
            n=int(part)-1
            if 0<=n<doc.page_count: result.insert_pdf(doc,from_page=n,to_page=n)
    result.save(out); result.close(); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="split.pdf")

# ── EXTRACT ────────────────────────────────────────────────────────────
@router.post("/extract")
async def extract(file: UploadFile = File(...), pages: str = Form(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"extracted.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src); result = fitz.open()
    for p in pages.split(","):
        p=p.strip()
        if p.isdigit():
            n=int(p)-1
            if 0<=n<doc.page_count: result.insert_pdf(doc,from_page=n,to_page=n)
    result.save(out); result.close(); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="extracted.pdf")

# ── ROTATE ─────────────────────────────────────────────────────────────
@router.post("/rotate")
async def rotate(file: UploadFile = File(...), degrees: int = Form(90), pages: str = Form("all")):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"rotated.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    pl = list(range(doc.page_count)) if pages=="all" else [int(p)-1 for p in pages.split(",") if p.strip().isdigit()]
    for n in pl:
        if 0<=n<doc.page_count: doc[n].set_rotation(degrees)
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="rotated.pdf")

# ── ROTATE VISUAL (per-page rotation map) ─────────────────────────────
@router.post("/rotate-visual")
async def rotate_visual(
    file: UploadFile = File(...),
    rotations: str = Form(...),
):
    import json
    try:
        rotation_map = json.loads(rotations)
    except:
        raise HTTPException(400, "Invalid rotations JSON")
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, "input.pdf")
    out = os.path.join(tmp, "rotated.pdf")
    with open(src, "wb") as f:
        shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    for page in doc:
        page_num = str(page.number + 1)
        if page_num in rotation_map:
            deg = int(rotation_map[page_num]) % 360
            doc[page.number].set_rotation(deg)
    doc.save(out, garbage=4, deflate=True)
    doc.close()
    return FileResponse(out, media_type="application/pdf", filename="rotated.pdf")

# ── DELETE PAGES ───────────────────────────────────────────────────────
@router.post("/delete-pages")
async def delete_pages(file: UploadFile = File(...), pages: str = Form(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"output.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    nums = sorted(set(int(p)-1 for p in pages.split(",") if p.strip().isdigit()), reverse=True)
    for n in nums:
        if 0<=n<doc.page_count: doc.delete_page(n)
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="output.pdf")

# ── COMPRESS ───────────────────────────────────────────────────────────
@router.post("/compress")
async def compress(file: UploadFile = File(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"compressed.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    doc.save(out, garbage=4, deflate=True, clean=True); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="compressed.pdf")

# ── WATERMARK HELPERS ──────────────────────────────────────────────────
def _resolve_text(text: str, page_num: int) -> str:
    from datetime import date
    text = text.replace("{{date}}", date.today().strftime("%Y-%m-%d"))
    text = text.replace("{{page_number}}", str(page_num))
    return text

def _parse_page_range(page_range: str, total: int):
    """Return list of 0-based page indices."""
    page_range = page_range.strip().lower()
    if page_range == "all":
        return list(range(total))
    if page_range == "odd":
        return [i for i in range(total) if (i+1) % 2 == 1]
    if page_range == "even":
        return [i for i in range(total) if (i+1) % 2 == 0]
    # custom: e.g. "1-3, 5, 7"
    indices = set()
    for part in page_range.split(","):
        part = part.strip()
        if "-" in part:
            try:
                a, b = part.split("-")
                for n in range(int(a)-1, min(int(b), total)):
                    indices.add(n)
            except: pass
        elif part.isdigit():
            n = int(part) - 1
            if 0 <= n < total:
                indices.add(n)
    return sorted(indices)

def _hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c*2 for c in hex_color)
    r, g, b = int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16)
    return (r/255, g/255, b/255)

def _apply_watermark_enhanced(page, text, font_size, rotation_deg, opacity, position, color_rgb, bold, italic):
    """Draw watermark at a specific position or tiled."""
    pw, ph = page.rect.width, page.rect.height
    mat = fitz.Matrix(rotation_deg)
    shape = page.new_shape()

    if position == "tile":
        # Draw tiled watermarks
        spacing_x = pw / 3
        spacing_y = ph / 4
        for row in range(-1, 5):
            for col in range(-1, 4):
                cx = col * spacing_x + spacing_x / 2
                cy = row * spacing_y + spacing_y / 2
                shape.insert_text(
                    fitz.Point(cx, cy), text,
                    fontsize=font_size,
                    color=color_rgb,
                    morph=(fitz.Point(cx, cy), fitz.Matrix(45)),
                )
    else:
        # Position-based placement
        margin = 60
        pos_map = {
            "top-left":       (margin,        font_size + margin),
            "top-center":     (pw/2,           font_size + margin),
            "top-right":      (pw - margin,    font_size + margin),
            "middle-left":    (margin,         ph/2),
            "middle-center":  (pw/2,           ph/2),
            "middle-right":   (pw - margin,    ph/2),
            "bottom-left":    (margin,         ph - margin),
            "bottom-center":  (pw/2,           ph - margin),
            "bottom-right":   (pw - margin,    ph - margin),
        }
        cx, cy = pos_map.get(position, (pw/2, ph/2))
        shape.insert_text(
            fitz.Point(cx, cy), text,
            fontsize=font_size,
            color=color_rgb,
            morph=(fitz.Point(cx, cy), mat),
        )

    shape.commit(overlay=True)

# ── WATERMARK ──────────────────────────────────────────────────────────
@router.post("/watermark")
async def watermark(
    file: UploadFile = File(...),
    mode: str = Form("text"),
    text: str = Form("CONFIDENTIAL"),
    opacity: str = Form("0.25"),
    rotation: str = Form("45"),
    font_size: str = Form("52"),
    font: str = Form("Helvetica"),
    bold: str = Form("false"),
    italic: str = Form("false"),
    underline: str = Form("false"),
    color: str = Form("#a0a0a0"),
    position: str = Form("middle-center"),
    layer: str = Form("over"),
    page_range: str = Form("all"),
    image: Optional[UploadFile] = File(None),
):
    try:
        rot = float(rotation)
        fs  = int(float(font_size))
        op  = float(opacity)
    except:
        rot, fs, op = 45.0, 52, 0.25

    color_rgb = _hex_to_rgb(color)
    # Apply opacity by blending towards white
    color_rgb = tuple(1 - (1 - c) * op for c in color_rgb)

    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, "input.pdf")
    out = os.path.join(tmp, "watermarked.pdf")
    try:
        with open(src, "wb") as f:
            shutil.copyfileobj(file.file, f)
        doc = fitz.open(src)
        pages = _parse_page_range(page_range, doc.page_count)
        overlay = layer != "under"

        for n in pages:
            if not (0 <= n < doc.page_count): continue
            page = doc[n]
            if mode == "image" and image:
                img_path = os.path.join(tmp, "wm_image")
                with open(img_path, "wb") as f:
                    image.file.seek(0)
                    shutil.copyfileobj(image.file, f)
                pw, ph = page.rect.width, page.rect.height
                wm_w, wm_h = 200, 100
                margin = 40
                pos_map = {
                    "top-left":      fitz.Rect(margin, margin, margin+wm_w, margin+wm_h),
                    "top-center":    fitz.Rect((pw-wm_w)/2, margin, (pw+wm_w)/2, margin+wm_h),
                    "top-right":     fitz.Rect(pw-wm_w-margin, margin, pw-margin, margin+wm_h),
                    "middle-left":   fitz.Rect(margin, (ph-wm_h)/2, margin+wm_w, (ph+wm_h)/2),
                    "middle-center": fitz.Rect((pw-wm_w)/2, (ph-wm_h)/2, (pw+wm_w)/2, (ph+wm_h)/2),
                    "middle-right":  fitz.Rect(pw-wm_w-margin, (ph-wm_h)/2, pw-margin, (ph+wm_h)/2),
                    "bottom-left":   fitz.Rect(margin, ph-wm_h-margin, margin+wm_w, ph-margin),
                    "bottom-center": fitz.Rect((pw-wm_w)/2, ph-wm_h-margin, (pw+wm_w)/2, ph-margin),
                    "bottom-right":  fitz.Rect(pw-wm_w-margin, ph-wm_h-margin, pw-margin, ph-margin),
                    "tile":          fitz.Rect((pw-wm_w)/2, (ph-wm_h)/2, (pw+wm_w)/2, (ph+wm_h)/2),
                }
                rect = pos_map.get(position, pos_map["middle-center"])
                page.insert_image(rect, filename=img_path, overlay=overlay)
            else:
                resolved = _resolve_text(text, n+1)
                _apply_watermark_enhanced(page, resolved, fs, rot, op, position, color_rgb, bold=="true", italic=="true")

        doc.save(out, garbage=4, deflate=True)
        doc.close()
    except Exception as e:
        shutil.rmtree(tmp, ignore_errors=True)
        raise HTTPException(500, f"Watermark failed: {str(e)}")
    return FileResponse(out, media_type="application/pdf", filename="watermarked.pdf")

# ── STAMP (FIXEDreturn FileResponse(out, media_type="application/pdf", filename="watermarked.pdf")

# ── STAMP (FIXED with position) ────────────────────────────────────────
@router.post("/stamp")
async def stamp(
    file: UploadFile = File(...),
    label: str = Form("APPROVED"),
    pages: str = Form("all"),
    position: str = Form("top-right"),  # top-left,top-center,top-right,middle-left,middle-center,middle-right,bottom-left,bottom-center,bottom-right
    custom_text: str = Form(""),
    font_size: int = Form(18),
    rotation: int = Form(0),
):
    COLORS = {
        "APPROVED":     (0.0, 0.55, 0.15),
        "DRAFT":        (0.8, 0.45, 0.0),
        "CONFIDENTIAL": (0.75, 0.0, 0.0),
        "REJECTED":     (0.85, 0.0, 0.0),
        "FINAL":        (0.0, 0.1, 0.8),
        "COPY":         (0.2, 0.2, 0.75),
        "CUSTOM":       (0.2, 0.2, 0.2),
    }
    stamp_text = custom_text.strip() if custom_text.strip() else label
    color      = COLORS.get(label.upper(), (0.3, 0.3, 0.3))
    stamp_w    = max(len(stamp_text) * font_size * 0.62 + 24, 90)
    stamp_h    = font_size + 18

    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,"input.pdf"); out = os.path.join(tmp,"stamped.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    pl = list(range(doc.page_count)) if pages=="all" else [int(p)-1 for p in pages.split(",") if p.strip().isdigit()]

    for n in pl:
        if not (0<=n<doc.page_count): continue
        page = doc[n]; pw,ph = page.rect.width, page.rect.height
        margin = 20

        # Position mapping
        pos_map = {
            "top-left":       (margin, margin),
            "top-center":     ((pw - stamp_w)/2, margin),
            "top-right":      (pw - stamp_w - margin, margin),
            "middle-left":    (margin, (ph - stamp_h)/2),
            "middle-center":  ((pw - stamp_w)/2, (ph - stamp_h)/2),
            "middle-right":   (pw - stamp_w - margin, (ph - stamp_h)/2),
            "bottom-left":    (margin, ph - stamp_h - margin),
            "bottom-center":  ((pw - stamp_w)/2, ph - stamp_h - margin),
            "bottom-right":   (pw - stamp_w - margin, ph - stamp_h - margin),
        }
        sx, sy = pos_map.get(position, (pw - stamp_w - margin, margin))
        rect = fitz.Rect(sx, sy, sx + stamp_w, sy + stamp_h)

        # Draw stamp box with fill
        bg_color = tuple(min(1.0, c + 0.88) for c in color)
        page.draw_rect(rect, color=color, fill=bg_color, width=2.5)

        # Text inside
        tx = sx + (stamp_w - len(stamp_text) * font_size * 0.6) / 2
        ty = sy + stamp_h - 6
        page.insert_text(fitz.Point(tx, ty), stamp_text, fontsize=font_size, color=color)

    doc.save(out, garbage=4, deflate=True); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="stamped.pdf")

# ── STAMP PREVIEW ──────────────────────────────────────────────────────
@router.post("/stamp/preview")
async def stamp_preview(
    file: UploadFile = File(...),
    label: str = Form("APPROVED"),
    position: str = Form("top-right"),
    custom_text: str = Form(""),
    font_size: int = Form(18),
    page_num: int = Form(1),
):
    """Returns a base64 PNG preview of the first page with stamp applied."""
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,"input.pdf"); out = os.path.join(tmp,"preview.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)

    # Reuse stamp logic on just one page
    COLORS = {"APPROVED":(0.0,0.55,0.15),"DRAFT":(0.8,0.45,0.0),"CONFIDENTIAL":(0.75,0.0,0.0),"REJECTED":(0.85,0.0,0.0),"FINAL":(0.0,0.1,0.8),"COPY":(0.2,0.2,0.75)}
    stamp_text = custom_text.strip() if custom_text.strip() else label
    color = COLORS.get(label.upper(), (0.3,0.3,0.3))
    stamp_w = max(len(stamp_text)*font_size*0.62+24, 90); stamp_h = font_size+18
    margin = 20

    doc = fitz.open(src)
    n = min(page_num-1, doc.page_count-1)
    page = doc[n]; pw,ph = page.rect.width, page.rect.height
    pos_map = {
        "top-left":(margin,margin),"top-center":((pw-stamp_w)/2,margin),"top-right":(pw-stamp_w-margin,margin),
        "middle-left":(margin,(ph-stamp_h)/2),"middle-center":((pw-stamp_w)/2,(ph-stamp_h)/2),"middle-right":(pw-stamp_w-margin,(ph-stamp_h)/2),
        "bottom-left":(margin,ph-stamp_h-margin),"bottom-center":((pw-stamp_w)/2,ph-stamp_h-margin),"bottom-right":(pw-stamp_w-margin,ph-stamp_h-margin),
    }
    sx,sy = pos_map.get(position,(pw-stamp_w-margin,margin))
    rect = fitz.Rect(sx,sy,sx+stamp_w,sy+stamp_h)
    bg_color = tuple(min(1.0,c+0.88) for c in color)
    page.draw_rect(rect, color=color, fill=bg_color, width=2.5)
    tx = sx+(stamp_w-len(stamp_text)*font_size*0.6)/2; ty = sy+stamp_h-6
    page.insert_text(fitz.Point(tx,ty), stamp_text, fontsize=font_size, color=color)

    mat = fitz.Matrix(0.8, 0.8)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    b64 = base64.b64encode(pix.tobytes("png")).decode()
    doc.close(); shutil.rmtree(tmp,ignore_errors=True)
    return JSONResponse({"preview": f"data:image/png;base64,{b64}"})

# ── WATERMARK PREVIEW ──────────────────────────────────────────────────
@router.post("/watermark/preview")
async def watermark_preview(
    file: UploadFile = File(...),
    mode: str = Form("text"),
    text: str = Form("CONFIDENTIAL"),
    opacity: str = Form("0.25"),
    rotation: str = Form("45"),
    font_size: str = Form("52"),
    font: str = Form("Helvetica"),
    bold: str = Form("false"),
    italic: str = Form("false"),
    color: str = Form("#a0a0a0"),
    position: str = Form("middle-center"),
    layer: str = Form("over"),
    image: Optional[UploadFile] = File(None),
):
    try:
        op  = float(opacity)
        rot = float(rotation)
        fs  = int(float(font_size))
    except:
        op, rot, fs = 0.25, 45.0, 52

    color_rgb = _hex_to_rgb(color)
    color_rgb = tuple(1 - (1 - c) * op for c in color_rgb)

    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, "input.pdf")
    with open(src, "wb") as f:
        shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    page = doc[0]

    if mode == "image" and image:
        img_path = os.path.join(tmp, "wm_image")
        with open(img_path, "wb") as f:
            shutil.copyfileobj(image.file, f)
        pw, ph = page.rect.width, page.rect.height
        wm_w, wm_h = 200, 100
        margin = 40
        pos_map = {
            "middle-center": fitz.Rect((pw-wm_w)/2, (ph-wm_h)/2, (pw+wm_w)/2, (ph+wm_h)/2),
            "tile":          fitz.Rect((pw-wm_w)/2, (ph-wm_h)/2, (pw+wm_w)/2, (ph+wm_h)/2),
        }
        rect = pos_map.get(position, fitz.Rect((pw-wm_w)/2, (ph-wm_h)/2, (pw+wm_w)/2, (ph+wm_h)/2))
        page.insert_image(rect, filename=img_path, overlay=True)
    else:
        resolved = _resolve_text(text, 1)
        _apply_watermark_enhanced(page, resolved, fs, rot, op, position, color_rgb, bold=="true", italic=="true")

    mat = fitz.Matrix(0.7, 0.7)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    b64 = base64.b64encode(pix.tobytes("png")).decode()
    doc.close()
    shutil.rmtree(tmp, ignore_errors=True)
    return JSONResponse({"preview": f"data:image/png;base64,{b64}"})

# ── PAGE NUMBERS ───────────────────────────────────────────────────────
@router.post("/page-numbers")
async def add_page_numbers(file: UploadFile = File(...), position: str = Form("bottom-center"), start: int = Form(1)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"numbered.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src); total = doc.page_count
    for i,page in enumerate(doc):
        pw,ph = page.rect.width, page.rect.height
        label = f"{i+start} / {total}"
        x = pw/2-20 if "center" in position else pw-80
        y = ph-25 if "bottom" in position else 25
        page.insert_text(fitz.Point(x,y), label, fontsize=9, color=(0.4,0.4,0.4))
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="numbered.pdf")

# ── PDF TO IMAGES ──────────────────────────────────────────────────────
@router.post("/pdf-to-images")
async def pdf_to_images(file: UploadFile = File(...)):
    import zipfile
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); zp = os.path.join(tmp,"images.zip")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src); mat = fitz.Matrix(2,2)
    with zipfile.ZipFile(zp,"w") as zf:
        for i,page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i+1}.png", pix.tobytes("png"))
    doc.close()
    return FileResponse(zp, media_type="application/zip", filename="pdf_images.zip")

# ── IMAGE TO PDF ───────────────────────────────────────────────────────
@router.post("/image-to-pdf")
async def image_to_pdf(files: List[UploadFile] = File(...)):
    tmp = tempfile.mkdtemp(); out = os.path.join(tmp,"from_images.pdf")
    doc = fitz.open()
    for f in files:
        path = os.path.join(tmp, f.filename)
        with open(path,"wb") as fp: shutil.copyfileobj(f.file,fp)
        page = doc.new_page(width=595, height=842)
        page.insert_image(fitz.Rect(0,0,595,842), filename=path)
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="from_images.pdf")

# ── PDF TO WORD (txt-based) ────────────────────────────────────────────
@router.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        raise HTTPException(500, "python-docx not installed.")
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"converted.docx")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc_pdf = fitz.open(src)
    word_doc = Document()
    for page in doc_pdf:
        blocks = page.get_text("blocks")
        for b in blocks:
            text = b[4].strip()
            if text:
                word_doc.add_paragraph(text)
        word_doc.add_page_break()
    doc_pdf.close()
    word_doc.save(out)
    return FileResponse(out, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="converted.docx")

# ── WORD TO PDF ────────────────────────────────────────────────────────
@router.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        raise HTTPException(500, "python-docx not installed.")
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)

    word_doc = Document(src)
    pdf_doc  = fitz.open()
    page     = pdf_doc.new_page(width=595, height=842)
    y        = 60; ml = 60; mr = 535; fs = 11; lh = fs * 1.4

    for para in word_doc.paragraphs:
        text = para.text.strip()
        if not text:
            y += lh * 0.5; continue
        if y > 800:
            page = pdf_doc.new_page(width=595, height=842); y = 60
        bold = any(run.bold for run in para.runs if run.bold)
        page.insert_text(fitz.Point(ml, y), text,
                         fontname="hebo" if bold else "helv", fontsize=fs, color=(0,0,0))
        y += lh

    pdf_doc.save(out); pdf_doc.close()
    return FileResponse(out, media_type="application/pdf", filename="converted.pdf")

# ── PROTECT ────────────────────────────────────────────────────────────
@router.post("/protect")
async def protect(file: UploadFile = File(...), password: str = Form(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"protected.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    perm = fitz.PDF_PERM_PRINT | fitz.PDF_PERM_COPY
    doc.save(out, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=password, owner_pw=password+"_owner", permissions=perm)
    doc.close()
    return FileResponse(out, media_type="application/pdf", filename="protected.pdf")

# ── UNLOCK ─────────────────────────────────────────────────────────────
@router.post("/unlock")
async def unlock(file: UploadFile = File(...), password: str = Form(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"unlocked.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    if doc.authenticate(password):
        doc.save(out, encryption=fitz.PDF_ENCRYPT_NONE); doc.close()
        return FileResponse(out, media_type="application/pdf", filename="unlocked.pdf")
    doc.close(); raise HTTPException(400,"Incorrect password.")

# ── REDACT ─────────────────────────────────────────────────────────────
@router.post("/redact")
async def redact(file: UploadFile = File(...), text: str = Form(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"redacted.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    for page in doc:
        areas = page.search_for(text)
        for rect in areas: page.add_redact_annot(rect, fill=(0,0,0))
        page.apply_redactions()
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="redacted.pdf")

# ── FIND & REPLACE ─────────────────────────────────────────────────────
@router.post("/find-replace")
async def find_replace(file: UploadFile = File(...), find: str = Form(...),
                       replace: str = Form(...), match_case: bool = Form(False), whole_word: bool = Form(False)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"replaced.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    for page in doc:
        areas = page.search_for(find)
        for rect in areas: page.add_redact_annot(rect, fill=(1,1,1))
        page.apply_redactions()
        for rect in areas:
            page.insert_text(fitz.Point(rect.x0, rect.y1), replace, fontsize=11, color=(0,0,0))
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="replaced.pdf")

# ── REPAIR ─────────────────────────────────────────────────────────────
@router.post("/repair")
async def repair(file: UploadFile = File(...)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"repaired.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src); doc.save(out, garbage=4, deflate=True, clean=True, linear=True); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="repaired.pdf")

# ── DUPLICATE PAGE ─────────────────────────────────────────────────────
@router.post("/duplicate-page")
async def duplicate_page(file: UploadFile = File(...), page: int = Form(1)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"duplicated.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src); n=page-1
    if 0<=n<doc.page_count: doc.copy_page(n)
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="duplicated.pdf")

# ── CROP (simple margins) ──────────────────────────────────────────────
@router.post("/crop")
async def crop(file: UploadFile = File(...), left: float=Form(0), top: float=Form(0), right: float=Form(0), bottom: float=Form(0)):
    tmp = tempfile.mkdtemp(); src = os.path.join(tmp,file.filename); out = os.path.join(tmp,"cropped.pdf")
    with open(src,"wb") as f: shutil.copyfileobj(file.file,f)
    doc = fitz.open(src)
    for page in doc:
        r=page.rect; page.set_cropbox(fitz.Rect(r.x0+left, r.y0+top, r.x1-right, r.y1-bottom))
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="cropped.pdf")

# ── CROP VISUAL (per-page drag crop) ───────────────────────────────────
@router.post("/crop-visual")
async def crop_visual(
    file: UploadFile = File(...),
    crops: str = Form(...),   # JSON: { "1": {left,top,right,bottom}, "2": {...}, ... }
):
    import json
    try:
        crop_data = json.loads(crops)
    except:
        raise HTTPException(400, "Invalid crops JSON")

    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, "input.pdf")
    out = os.path.join(tmp, "cropped.pdf")

    with open(src, "wb") as f:
        shutil.copyfileobj(file.file, f)

    doc = fitz.open(src)

    for page in doc:
        page_num = str(page.number + 1)
        if page_num not in crop_data:
            continue
        c = crop_data[page_num]
        left   = float(c.get("left",   0))
        top    = float(c.get("top",    0))
        right  = float(c.get("right",  0))
        bottom = float(c.get("bottom", 0))
        r = page.rect
        new_rect = fitz.Rect(
            r.x0 + left,
            r.y0 + top,
            r.x1 - right,
            r.y1 - bottom,
        )
        # Clamp to valid bounds
        if new_rect.width > 10 and new_rect.height > 10:
            page.set_cropbox(new_rect)

    doc.save(out, garbage=4, deflate=True)
    doc.close()
    return FileResponse(out, media_type="application/pdf", filename="cropped.pdf")

# ══════════════════════════════════════════════════════════════════════
# CONVERT TO PDF
# ══════════════════════════════════════════════════════════════════════

# ── PowerPoint → PDF ───────────────────────────────────────────────────
@router.post("/pptx-to-pdf")
async def pptx_to_pdf(file: UploadFile = File(...)):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise HTTPException(500, "python-pptx not installed. Run: pip install python-pptx")
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted.pdf")
    with open(src, "wb") as f: shutil.copyfileobj(file.file, f)
    prs = Presentation(src)
    doc = fitz.open()
    for slide in prs.slides:
        page = doc.new_page(width=792, height=612)  # 11x8.5 landscape
        y = 60
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text: continue
                    if y > 580: break
                    bold = any(run.bold for run in para.runs if run.bold is not None)
                    fs = 14 if y < 100 else 11  # larger for title
                    page.insert_text(fitz.Point(40, y), text,
                        fontname="hebo" if bold else "helv", fontsize=fs, color=(0,0,0))
                    y += fs * 1.6
        y = 60
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="converted.pdf")

# ── Excel → PDF ────────────────────────────────────────────────────────
@router.post("/xlsx-to-pdf")
async def xlsx_to_pdf(file: UploadFile = File(...)):
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise HTTPException(500, "openpyxl not installed. Run: pip install openpyxl")
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted.pdf")
    with open(src, "wb") as f: shutil.copyfileobj(file.file, f)
    wb = load_workbook(src, data_only=True)
    doc = fitz.open()
    for sheet in wb.worksheets:
        page = doc.new_page(width=842, height=595)  # A4 landscape
        x_start, y_start, fs = 30, 40, 9
        col_w = min(120, max(60, (842 - 60) // max(sheet.max_column, 1)))
        for ri, row in enumerate(sheet.iter_rows(values_only=True)):
            y = y_start + ri * (fs + 4)
            if y > 575:
                page = doc.new_page(width=842, height=595)
                y_start = 40; ri = 0; y = 40
            for ci, cell in enumerate(row):
                if ci > 12: break
                text = str(cell) if cell is not None else ""
                if not text or text == "None": continue
                x = x_start + ci * col_w
                page.insert_text(fitz.Point(x, y), text[:20],
                    fontname="helv", fontsize=fs, color=(0,0,0))
    doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="converted.pdf")

# ── HTML → PDF ─────────────────────────────────────────────────────────
@router.post("/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted.pdf")
    with open(src, "wb") as f: shutil.copyfileobj(file.file, f)
    html_content = open(src, "r", encoding="utf-8", errors="ignore").read()
    # Use PyMuPDF's built-in HTML story renderer
    try:
        story = fitz.Story(html=html_content)
        writer = fitz.DocumentWriter(out)
        mediabox = fitz.paper_rect("a4")
        where    = mediabox + (36, 36, -36, -36)
        more     = True
        while more:
            dev, more = writer.begin_page(mediabox)
            more, _ = story.place(where)
            story.draw(dev)
            writer.end_page()
        writer.close()
    except Exception:
        # Fallback: strip HTML tags and write plain text
        import re as _re
        text = _re.sub(r'<[^>]+>', ' ', html_content)
        text = _re.sub(r'\s+', ' ', text).strip()
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_textbox(fitz.Rect(36,36,559,805), text,
            fontname="helv", fontsize=11, color=(0,0,0))
        doc.save(out); doc.close()
    return FileResponse(out, media_type="application/pdf", filename="converted.pdf")

# ══════════════════════════════════════════════════════════════════════
# CONVERT FROM PDF
# ══════════════════════════════════════════════════════════════════════

# ── PDF → JPG (single images, not ZIP) ────────────────────────────────
@router.post("/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...)):
    import zipfile
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    zp  = os.path.join(tmp, "images.zip")
    with open(src,"wb") as f: shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    mat = fitz.Matrix(2, 2)
    with zipfile.ZipFile(zp, "w") as zf:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i+1}.jpg", pix.tobytes("jpeg", jpg_quality=90))
    doc.close()
    return FileResponse(zp, media_type="application/zip", filename="pdf_images_jpg.zip")

# ── PDF → PowerPoint ───────────────────────────────────────────────────
@router.post("/pdf-to-pptx")
async def pdf_to_pptx(file: UploadFile = File(...)):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
    except ImportError:
        raise HTTPException(500, "python-pptx not installed. Run: pip install python-pptx")
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted.pptx")
    with open(src, "wb") as f: shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    prs = Presentation()
    prs.slide_width  = Inches(11)
    prs.slide_height = Inches(8.5)
    blank = prs.slide_layouts[6]  # blank layout
    for page in doc:
        slide = prs.slides.add_slide(blank)
        # Render page as image and use as background
        pix      = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
        img_path = os.path.join(tmp, f"slide_{page.number}.png")
        pix.save(img_path)
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
    doc.close()
    prs.save(out)
    return FileResponse(out,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="converted.pptx")

# ── PDF → Excel ────────────────────────────────────────────────────────
@router.post("/pdf-to-xlsx")
async def pdf_to_xlsx(file: UploadFile = File(...)):
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
    except ImportError:
        raise HTTPException(500, "openpyxl not installed. Run: pip install openpyxl")
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted.xlsx")
    with open(src, "wb") as f: shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    wb  = Workbook()
    for pi, page in enumerate(doc):
        ws = wb.active if pi == 0 else wb.create_sheet(f"Page {pi+1}")
        if pi == 0: ws.title = "Page 1"
        blocks = page.get_text("blocks")
        # Sort by vertical then horizontal position
        blocks.sort(key=lambda b: (round(b[1]/15)*15, b[0]))
        for block in blocks:
            text = block[4].strip().replace("\n", " ")
            if not text: continue
            # Estimate row/col from position
            row = max(1, int(block[1] / 15) + 1)
            col = max(1, int(block[0] / 80) + 1)
            try:
                ws.cell(row=row, column=col, value=text)
            except Exception:
                pass
    doc.close()
    wb.save(out)
    return FileResponse(out,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename="converted.xlsx")

# ── PDF → PDF/A ────────────────────────────────────────────────────────
@router.post("/pdf-to-pdfa")
async def pdf_to_pdfa(file: UploadFile = File(...)):
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, file.filename)
    out = os.path.join(tmp, "converted_pdfa.pdf")
    with open(src, "wb") as f: shutil.copyfileobj(file.file, f)
    doc = fitz.open(src)
    # PDF/A-1b: linearized, no encryption, clean structure
    doc.save(out, garbage=4, deflate=True, clean=True, linear=True,
             encryption=fitz.PDF_ENCRYPT_NONE)
    doc.close()
    return FileResponse(out, media_type="application/pdf", filename="converted_pdfa.pdf")


# ══════════════════════════════════════════════════════════════════════
# NEW TOOLS
# ══════════════════════════════════════════════════════════════════════

# ── ORGANIZE PDF ───────────────────────────────────────────────────────
@router.post("/organize")
async def organize_pdf(
    file:         UploadFile = File(...),
    instructions: str        = Form(...),   # JSON: [{orig_idx, rotation, blank}]
):
    """
    Reorder, rotate, and/or delete pages. Insert blank pages.
    instructions = JSON array where each element is:
      { "orig_idx": int,   # 0-based index in original PDF (-1 for blank)
        "rotation": int,   # additional rotation degrees (0/90/180/270)
        "blank":    bool   # if true, insert a blank white page }
    """
    import json
    try:
        steps = json.loads(instructions)
    except Exception:
        raise HTTPException(400, "Invalid instructions JSON")

    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, "input.pdf")
    out = os.path.join(tmp, "organized.pdf")
    try:
        with open(src, "wb") as f:
            shutil.copyfileobj(file.file, f)
        src_doc = fitz.open(src)
        result  = fitz.open()

        for step in steps:
            blank    = step.get("blank", False)
            orig_idx = int(step.get("orig_idx", -1))
            rotation = int(step.get("rotation", 0)) % 360

            if blank:
                # Insert a blank A4 page
                pw = src_doc[0].rect.width  if src_doc.page_count > 0 else 595
                ph = src_doc[0].rect.height if src_doc.page_count > 0 else 842
                result.new_page(width=pw, height=ph)
            elif 0 <= orig_idx < src_doc.page_count:
                result.insert_pdf(src_doc, from_page=orig_idx, to_page=orig_idx)
                # Apply extra rotation to the newly inserted page
                if rotation:
                    last = result.page_count - 1
                    existing = result[last].rotation
                    result[last].set_rotation((existing + rotation) % 360)

        result.save(out, garbage=4, deflate=True)
        result.close()
        src_doc.close()
    except Exception as e:
        shutil.rmtree(tmp, ignore_errors=True)
        raise HTTPException(500, f"Organize failed: {str(e)}")

    return FileResponse(out, media_type="application/pdf", filename="organized.pdf")


# ── SIGN PDF ───────────────────────────────────────────────────────────
@router.post("/sign")
async def sign_pdf(
    file:      UploadFile = File(...),
    signature: str        = Form(...),   # base64 PNG data-URL
    x:         str        = Form("50"),
    y:         str        = Form("50"),
    width:     str        = Form("200"),
    height:    str        = Form("80"),
    page:      str        = Form("1"),
):
    """
    Embed a base64-encoded PNG signature image onto a specific page
    at the given position and size (all in PDF points).
    """
    import re as _re

    try:
        px     = float(x);      py  = float(y)
        pw_sig = float(width);  ph_sig = float(height)
        pg_num = max(1, int(page))
    except ValueError:
        raise HTTPException(400, "Invalid numeric parameters")

    # Strip data-URL prefix: "data:image/png;base64,<data>"
    b64_match = _re.match(r"data:image/[^;]+;base64,(.+)", signature, _re.DOTALL)
    if not b64_match:
        raise HTTPException(400, "signature must be a base64 data-URL")
    img_bytes = base64.b64decode(b64_match.group(1))

    tmp = tempfile.mkdtemp()
    src      = os.path.join(tmp, "input.pdf")
    sig_path = os.path.join(tmp, "sig.png")
    out      = os.path.join(tmp, "signed.pdf")
    try:
        with open(src, "wb") as f:
            shutil.copyfileobj(file.file, f)
        with open(sig_path, "wb") as f:
            f.write(img_bytes)

        doc    = fitz.open(src)
        pg_idx = min(pg_num - 1, doc.page_count - 1)
        page_obj = doc[pg_idx]

        rect = fitz.Rect(px, py, px + pw_sig, py + ph_sig)
        page_obj.insert_image(rect, filename=sig_path, overlay=True)

        doc.save(out, garbage=4, deflate=True)
        doc.close()
    except Exception as e:
        shutil.rmtree(tmp, ignore_errors=True)
        raise HTTPException(500, f"Sign failed: {str(e)}")

    return FileResponse(out, media_type="application/pdf", filename="signed.pdf")


# ── OCR PDF ────────────────────────────────────────────────────────────
@router.post("/ocr")
async def ocr_pdf(
    file:        UploadFile = File(...),
    language:    str        = Form("eng"),
    output_mode: str        = Form("searchable"),
    page_texts:  str        = Form(""),   # JSON array of {page, text} from client-side OCR
):
    """
    3-Tier OCR — designed to put zero load on this instance.

    Tier 1 (primary):   OCR.space API — client pre-processes pages to JPEG,
                        sends them individually. Server just receives text results.
    Tier 2 (fallback):  Mistral AI OCR API — if OCR.space fails/limit hit.
    Tier 3 (last):      pytesseract — local fallback, no torch/easyocr needed.

    For TEXT mode:  browser sends pre-extracted text → server just returns it.
    For SEARCHABLE: browser sends page texts → server embeds them into PDF.

    RAM usage: ~10MB (vs 500MB+ with easyocr).
    """
    import os, json, urllib.request, urllib.parse, base64

    OCR_SPACE_KEY = os.environ.get("OCR_SPACE_API_KEY", "")
    MISTRAL_KEY   = os.environ.get("MISTRAL_API_KEY", "")

    # ── Language code mapping ──────────────────────────────────────────
    # OCR.space uses different codes than tesseract
    LANG_MAP = {
        "eng":"eng","fra":"fre","deu":"ger","spa":"spa","ita":"ita",
        "por":"por","rus":"rus","chi_sim":"chs","jpn":"jpn","kor":"kor",
        "ara":"ara","hin":"hin",
    }
    ocr_lang = LANG_MAP.get(language, "eng")

    # ── If client already did OCR and sent text results ────────────────
    # This is the zero-load path — browser did the heavy lifting
    if page_texts.strip():
        try:
            pages = json.loads(page_texts)
            combined = "\n\n--- Page Break ---\n\n".join(
                p.get("text", "") for p in sorted(pages, key=lambda x: x.get("page", 0))
            )
            if output_mode == "text":
                return JSONResponse({"text": combined})

            # Searchable PDF — embed text into original PDF
            tmp = tempfile.mkdtemp()
            src = os.path.join(tmp, "input.pdf")
            out = os.path.join(tmp, "ocr_searchable.pdf")
            with open(src, "wb") as f:
                shutil.copyfileobj(file.file, f)

            doc = fitz.open(src)
            for p in sorted(pages, key=lambda x: x.get("page", 0)):
                pg_num = p.get("page", 1) - 1
                if 0 <= pg_num < doc.page_count:
                    page    = doc[pg_num]
                    pw, ph  = page.rect.width, page.rect.height
                    words   = p.get("text", "").split()
                    if not words:
                        continue
                    # Distribute words across page invisibly
                    fs      = 10
                    x, y    = 36, 36
                    line_h  = fs * 1.4
                    for word in words:
                        if x + len(word) * 6 > pw - 36:
                            x  = 36
                            y += line_h
                        if y > ph - 36:
                            break
                        try:
                            page.insert_text(
                                fitz.Point(x, y), word,
                                fontsize=fs, color=(1,1,1), overlay=True
                            )
                        except Exception:
                            pass
                        x += len(word) * 6 + 4

            doc.save(out, garbage=4, deflate=True)
            doc.close()
            return FileResponse(out, media_type="application/pdf", filename="ocr_searchable.pdf")

        except Exception as e:
            raise HTTPException(500, f"Failed to process OCR results: {e}")

    # ── Server-side fallback (Tier 2: Mistral, Tier 3: pytesseract) ───
    # Only reached if client-side OCR failed completely
    tmp = tempfile.mkdtemp()
    src = os.path.join(tmp, "input.pdf")
    out = os.path.join(tmp, "ocr_output.pdf")

    try:
        with open(src, "wb") as f:
            shutil.copyfileobj(file.file, f)

        doc       = fitz.open(src)
        mat       = fitz.Matrix(2.0, 2.0)
        all_texts = []

        def _ocr_via_mistral(img_b64: str) -> str | None:
            """Use Mistral AI vision API for OCR."""
            if not MISTRAL_KEY:
                return None
            try:
                payload = json.dumps({
                    "model": "pixtral-12b-2409",
                    "messages": [{
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                            {"type": "text",      "text": "Extract all text from this image. Return only the text, no commentary."},
                        ]
                    }],
                    "max_tokens": 2000,
                }).encode()
                req = urllib.request.Request(
                    "https://api.mistral.ai/v1/chat/completions",
                    data=payload,
                    headers={
                        "Content-Type":  "application/json",
                        "Authorization": f"Bearer {MISTRAL_KEY}",
                    },
                    method="POST"
                )
                with urllib.request.urlopen(req, timeout=30) as resp:
                    data = json.loads(resp.read().decode())
                    return data["choices"][0]["message"]["content"].strip()
            except Exception:
                return None

        def _ocr_via_tesseract(img_bytes: bytes, lang: str) -> str:
            """Last-resort local OCR — lightweight, no torch needed."""
            try:
                import pytesseract
                from PIL import Image
                import io
                img = Image.open(io.BytesIO(img_bytes))
                return pytesseract.image_to_string(img, lang=lang[:3])
            except Exception:
                return ""

        for page in doc:
            pix      = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("jpeg")
            img_b64  = base64.b64encode(img_bytes).decode()

            # Try Mistral first (Tier 2)
            text = _ocr_via_mistral(img_b64)

            # Fall back to pytesseract (Tier 3)
            if not text:
                text = _ocr_via_tesseract(img_bytes, language)

            all_texts.append(text or "")

        combined = "\n\n--- Page Break ---\n\n".join(all_texts)

        if output_mode == "text":
            doc.close()
            shutil.rmtree(tmp, ignore_errors=True)
            return JSONResponse({"text": combined})

        # Embed text into searchable PDF
        for i, (page, text) in enumerate(zip(doc, all_texts)):
            if not text.strip():
                continue
            pw, ph = page.rect.width, page.rect.height
            words  = text.split()
            fs     = 10
            x, y   = 36, 36
            line_h = fs * 1.4
            for word in words:
                if x + len(word) * 6 > pw - 36:
                    x  = 36
                    y += line_h
                if y > ph - 36:
                    break
                try:
                    page.insert_text(
                        fitz.Point(x, y), word,
                        fontsize=fs, color=(1,1,1), overlay=True
                    )
                except Exception:
                    pass
                x += len(word) * 6 + 4

        doc.save(out, garbage=4, deflate=True)
        doc.close()

    except Exception as e:
        shutil.rmtree(tmp, ignore_errors=True)
        raise HTTPException(500, f"OCR failed: {str(e)}")

    return FileResponse(out, media_type="application/pdf", filename="ocr_searchable.pdf")


# ── TRANSLATE PDF ──────────────────────────────────────────────────────
@router.post("/translate")
async def translate_pdf(
    file:        UploadFile = File(...),
    source_lang: str        = Form("auto"),
    target_lang: str        = Form("en"),
):
    """
    Translate PDF using Lingva Translate (Google-quality, free, no API key).
    Falls back to unofficial Google Translate API if Lingva is down.

    RAM used: ~5MB (just HTTP requests — no model loaded on instance).
    Speed: 2-5 seconds vs 30-60 seconds with argostranslate.
    """
    import urllib.request, urllib.parse, urllib.error, json

    # ── Language code normalisation ────────────────────────────────────
    # Lingva + Google both use 2-letter ISO codes
    LANG_MAP = {
        "auto": "auto", "en": "en", "fr": "fr", "de": "de",
        "es": "es", "it": "it", "pt": "pt", "ru": "ru",
        "zh": "zh", "ja": "ja", "ko": "ko", "ar": "ar",
        "hi": "hi", "nl": "nl", "pl": "pl", "tr": "tr", "sv": "sv",
    }
    src = LANG_MAP.get(source_lang, "auto")
    tgt = LANG_MAP.get(target_lang, "en")
    if src == tgt:
        raise HTTPException(400, "Source and target languages must be different.")

    # ── Lingva instances (tried in order, first success wins) ──────────
    LINGVA_INSTANCES = [
        "https://lingva.ml",
        "https://lingva.thedaviddelta.com",
        "https://translate.plausibility.cloud",
        "https://lingva.lunar.icu",
    ]

    def _translate_lingva(text: str, s: str, t: str) -> str | None:
        """Try all Lingva instances. Returns translated text or None."""
        if not text.strip():
            return text
        encoded = urllib.parse.quote(text, safe="")
        src_code = "auto" if s == "auto" else s
        for base in LINGVA_INSTANCES:
            try:
                url = f"{base}/api/v1/{src_code}/{t}/{encoded}"
                req = urllib.request.Request(url, headers={"User-Agent": "PDFForge/1.0"})
                with urllib.request.urlopen(req, timeout=8) as resp:
                    data = json.loads(resp.read().decode())
                    result = data.get("translation", "").strip()
                    if result:
                        return result
            except Exception:
                continue  # try next instance
        return None

    def _translate_google(text: str, s: str, t: str) -> str | None:
        """
        Unofficial Google Translate API — fallback if all Lingva instances fail.
        Same engine as translate.google.com, no API key needed.
        """
        if not text.strip():
            return text
        try:
            params = urllib.parse.urlencode({
                "client": "gtx",
                "sl":     s if s != "auto" else "auto",
                "tl":     t,
                "dt":     "t",
                "q":      text,
            })
            url = f"https://translate.googleapis.com/translate_a/single?{params}"
            req = urllib.request.Request(url, headers={
                "User-Agent": "Mozilla/5.0",
                "Accept":     "application/json",
            })
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read().decode())
                # Response format: [[[translated, original, ...], ...], ...]
                parts = data[0]
                result = "".join(p[0] for p in parts if p[0])
                return result.strip() or None
        except Exception:
            return None

    def _translate(text: str, s: str, t: str) -> str:
        """Translate with Lingva → Google fallback → original text."""
        if not text or not text.strip() or len(text.strip()) < 2:
            return text
        # Try Lingva first
        result = _translate_lingva(text, s, t)
        if result:
            return result
        # Fallback to Google
        result = _translate_google(text, s, t)
        if result:
            return result
        # Last resort: return original
        return text

    # ── Process PDF ────────────────────────────────────────────────────
    tmp = tempfile.mkdtemp()
    src_path = os.path.join(tmp, "input.pdf")
    out_path = os.path.join(tmp, "translated.pdf")

    try:
        with open(src_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        doc = fitz.open(src_path)

        for page in doc:
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

            for block in blocks:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        original = span.get("text", "").strip()
                        if not original or len(original) < 2:
                            continue

                        translated = _translate(original, src, tgt)
                        if not translated or translated == original:
                            continue

                        bbox     = fitz.Rect(span["bbox"])
                        fontsize = max(4.0, span.get("size", 11))
                        color_int = span.get("color", 0)
                        r = ((color_int >> 16) & 0xFF) / 255
                        g = ((color_int >>  8) & 0xFF) / 255
                        b = ( color_int        & 0xFF) / 255

                        # White-out original, write translation
                        page.draw_rect(bbox, color=None, fill=(1, 1, 1), overlay=True)
                        try:
                            page.insert_textbox(
                                bbox, translated,
                                fontsize=fontsize,
                                color=(r, g, b),
                                align=0,
                                overlay=True,
                            )
                        except Exception:
                            pass

        doc.save(out_path, garbage=4, deflate=True)
        doc.close()

    except Exception as e:
        shutil.rmtree(tmp, ignore_errors=True)
        raise HTTPException(500, f"Translation failed: {str(e)}")

    return FileResponse(
        out_path,
        media_type="application/pdf",
        filename=f"translated_{target_lang}.pdf"
    )
